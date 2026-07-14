"""폴더 안의 PDF/xlsx/pptx/txt 파일을 읽어 LangChain Document 리스트로 변환."""
import os
from typing import List

import fitz  # PyMuPDF
import pandas as pd
from pptx import Presentation
from langchain_core.documents import Document


def load_pdf(path: str) -> List[Document]:
    docs = []
    with fitz.open(path) as pdf:
        for i, page in enumerate(pdf):
            text = page.get_text().strip()
            if text:
                docs.append(Document(
                    page_content=text,
                    metadata={"source": os.path.basename(path), "page": i + 1, "type": "pdf"},
                ))
    return docs


def load_xlsx(path: str) -> List[Document]:
    docs = []
    xls = pd.ExcelFile(path)
    for sheet in xls.sheet_names:
        df = xls.parse(sheet).fillna("")
        if df.empty:
            continue
        # 시트를 마크다운 표 문자열로 변환 (LLM이 표를 잘 이해)
        text = df.to_markdown(index=False)
        docs.append(Document(
            page_content=f"[시트: {sheet}]\n{text}",
            metadata={"source": os.path.basename(path), "sheet": sheet, "type": "xlsx"},
        ))
    return docs


def load_pptx(path: str) -> List[Document]:
    docs = []
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
        text = "\n".join(p for p in parts if p.strip())
        if text.strip():
            docs.append(Document(
                page_content=text,
                metadata={"source": os.path.basename(path), "slide": i + 1, "type": "pptx"},
            ))
    return docs


def load_txt(path: str) -> List[Document]:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read().strip()
    if not text:
        return []
    return [Document(
        page_content=text,
        metadata={"source": os.path.basename(path), "type": "txt"},
    )]


LOADERS = {
    ".pdf": load_pdf,
    ".xlsx": load_xlsx,
    ".xls": load_xlsx,
    ".pptx": load_pptx,
    ".txt": load_txt,
    ".md": load_txt,
}


def load_directory(docs_dir: str) -> List[Document]:
    """폴더를 재귀 순회하며 지원 형식 파일을 모두 로드."""
    all_docs: List[Document] = []
    for root, _, files in os.walk(docs_dir):
        for name in files:
            ext = os.path.splitext(name)[1].lower()
            loader = LOADERS.get(ext)
            if not loader:
                continue
            path = os.path.join(root, name)
            try:
                all_docs.extend(loader(path))
                print(f"  loaded: {name}")
            except Exception as e:
                print(f"  [skip] {name}: {e}")
    return all_docs
